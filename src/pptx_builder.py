import os
import logging
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger("DAIS_Intraday.PPTXBuilder")

def build_presentation(summary_list: list, cfg: dict, outdir: str) -> str:
    """
    Constructs an enterprise-grade PowerPoint presentation summarizing high-frequency metrics.
    Ensures safe shape placement grids and prevents image path corruption crashes.
    """
    pptx_filename = os.path.join(outdir, "DAIS_Performance_Executive_Deck.pptx")
    
    prs = Presentation()
    # Force standard 16:9 widescreen presentation dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layout templates to construct custom layout structures manually
    blank_layout = prs.slide_layouts[6]
    
    # Custom high-contrast executive design theme palette
    COLOR_DARK = RGBColor(30, 41, 59)      # Slate Blue Background
    COLOR_LIGHT = RGBColor(248, 250, 252)  # Light Canvas Background
    COLOR_SKY = RGBColor(2, 132, 199)      # Accent Blue Text
    COLOR_WHITE = RGBColor(255, 255, 255)
    
    # ------------------------------------------------------------
    # SLIDE 1: COVER/TITLE DECK TITLE
    # ------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    
    # Construct corporate background panel
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLOR_DARK
    bg_shape.line.fill.background() # Clear borders
    
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.333), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "DAIS Project: High-Frequency Strategy Results"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = f"Intraday Engine Telemetry Brief | Horizon: {cfg.get('lookback_days', 30)} Days | Interval: {cfg.get('bar_interval', '5m')}"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_SKY
    p2.space_before = Pt(15)

    # ------------------------------------------------------------
    # SLIDE 2: GLOBAL PERFORMANCE OVERVIEW TABLE MATRIX
    # ------------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Global Portfolio Performance Matrix"
    title_box.text_frame.paragraphs[0].font.size = Pt(22)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK
    
    rows = len(summary_list) + 1
    cols = 6
    left = Inches(0.75)
    top = Inches(1.5)
    width = Inches(11.833)
    height = Inches(0.4 * rows)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["Ticker", "Total Return", "Ann. Return", "Sharpe Ratio", "Max Drawdown", "True Beta"]
    col_widths = [Inches(1.833), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = w
        
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE

    for row_idx, item in enumerate(summary_list, start=1):
        m = item.get("metrics", {})
        
        def fmt_pct(val):
            return f"{val * 100:.2f}%" if (val is not None and not np.isnan(val)) else "0.00%"
        def fmt_num(val):
            return f"{val:.2f}" if (val is not None and not np.isnan(val)) else "N/A"
            
        row_data = [
            item.get("ticker", "UNKNOWN"),
            fmt_pct(m.get("total_return")),
            fmt_pct(m.get("annualized_return")),
            fmt_num(m.get("sharpe")),
            fmt_pct(m.get("max_drawdown")),
            fmt_num(item.get("beta_true"))
        ]
        
        for col_idx, val_str in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = val_str
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(11)
                p.font.color.rgb = COLOR_DARK

    # ------------------------------------------------------------
    # SLIDES 3+: ASSET LEVEL DETAIL VISUALIZATIONS
    # ------------------------------------------------------------
    for item in summary_list:
        if "error" in item:
            continue
            
        ticker = item.get("ticker", "UNKNOWN")
        charts = item.get("charts", {})
        sell_stats = item.get("sell_rule_stats", {})
        
        slide = prs.slides.add_slide(blank_layout)
        
        # Header title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(0.6))
        p = title_box.text_frame.paragraphs[0]
        p.text = f"Execution Performance Stream: {ticker}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_SKY
        
        # Detailed stats subtext banner box
        stats_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(11.833), Inches(0.5))
        p_stats = stats_box.text_frame.paragraphs[0]
        p_stats.text = (
            f"Executed Sells: {sell_stats.get('executed_sells', 0)} | "
            f"Overnight Liquidations: {sell_stats.get('overnight_liquidations', 0)} | "
            f"Rejected via MA20: {sell_stats.get('rejected_ma20', 0)} | "
            f"Rejected via Cost Basis: {sell_stats.get('rejected_avg_cost', 0)}"
        )
        p_stats.font.size = Pt(11)
        p_stats.font.color.rgb = RGBColor(100, 116, 139)
        
        # Build strict non-overlapping image layout grids
        topo_img = charts.get("topology")
        equity_img = charts.get("equity_curve")
        inventory_img = charts.get("inventory_density")
        
        # Compute dynamic horizontal placement coordinate markers
        # Slot 1: Topology Plot
        if topo_img and os.path.exists(topo_img):
            slide.shapes.add_picture(topo_img, Inches(0.5), Inches(1.8), width=Inches(6.0), height=Inches(2.5))
            lbl = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(6.0), Inches(0.3))
            lbl.text_frame.paragraphs[0].text = "Execution Strategy Topology Profiles"
            lbl.text_frame.paragraphs[0].font.size = Pt(10)
            lbl.text_frame.paragraphs[0].font.italic = True
            
        # Slot 2: Capital Balance Growth Curve
        if equity_img and os.path.exists(equity_img):
            slide.shapes.add_picture(equity_img, Inches(6.833), Inches(1.8), width=Inches(6.0), height=Inches(2.5))
            lbl = slide.shapes.add_textbox(Inches(6.833), Inches(4.3), Inches(6.0), Inches(0.3))
            lbl.text_frame.paragraphs[0].text = "Accumulated Growth Value Optimization"
            lbl.text_frame.paragraphs[0].font.size = Pt(10)
            lbl.text_frame.paragraphs[0].font.italic = True
            
        # Slot 3: Large Footprint Consolidated Inventory Density Plot (Spans Bottom Margin width)
        if inventory_img and os.path.exists(inventory_img):
            slide.shapes.add_picture(inventory_img, Inches(0.5), Inches(4.7), width=Inches(12.333), height=Inches(2.1))
            lbl = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.3))
            lbl.text_frame.paragraphs[0].text = "Intraday Base Share Inventory Density Allocations"
            lbl.text_frame.paragraphs[0].font.size = Pt(10)
            lbl.text_frame.paragraphs[0].font.italic = True

    try:
        prs.save(pptx_filename)
        logger.info(f"High-frequency PowerPoint brief deck exported cleanly to: {pptx_filename}")
        return pptx_filename
    except Exception as e:
        logger.error(f"Failed to generate layout shapes inside python-pptx: {e}", exc_info=True)
        raise e
